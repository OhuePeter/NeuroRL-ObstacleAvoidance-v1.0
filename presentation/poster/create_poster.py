"""
=============================================================
Neural Population Dynamics Poster Generator

Authors
-------
Peter Ohue
Emily Oby
Gunnar Blohm

Queen's University

Generates a fully editable PowerPoint conference poster.

=============================================================
"""

from pathlib import Path
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Optional dependencies for QR generation and PDF export.
try:
    import qrcode
except Exception:
    qrcode = None

try:
    import win32com.client  # type: ignore
except Exception:
    win32com = None

# ----------------------------------------------------------
# Directories
# ----------------------------------------------------------

ROOT = Path(__file__).parent

FIGURES = ROOT / "figures"

LOGOS = ROOT / "logos"

OUTPUT = ROOT / "output"

OUTPUT.mkdir(exist_ok=True)

POSTER_QR_URL = "https://github.com/OhuePeter/NeuroRL-ObstacleAvoidance-v1.0"

# ----------------------------------------------------------
# Poster
# ----------------------------------------------------------

prs = Presentation()

prs.slide_width = Inches(48)

prs.slide_height = Inches(27)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ----------------------------------------------------------
# Colours
# ----------------------------------------------------------

QUEENS_BLUE = RGBColor(0, 53, 95)

LIGHT_BLUE = RGBColor(224, 238, 255)

DARK = RGBColor(30, 30, 30)

GRAY = RGBColor(110, 110, 110)

WHITE = RGBColor(255,255,255)

# ----------------------------------------------------------
# Helpers
# ----------------------------------------------------------

def textbox(x,y,w,h,text,size=18,bold=False,color=DARK):

    tb = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    tf = tb.text_frame

    p = tf.paragraphs[0]

    p.text = text

    p.font.size = Pt(size)

    p.font.bold = bold

    p.font.color.rgb = color

    return tb

def heading(title,x,y):

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(14),
        Inches(0.55)
    )

    shape.fill.solid()

    shape.fill.fore_color.rgb = QUEENS_BLUE

    shape.line.color.rgb = QUEENS_BLUE

    tf = shape.text_frame

    p = tf.paragraphs[0]

    p.text = title

    p.font.size = Pt(22)

    p.font.bold = True

    p.font.color.rgb = WHITE

def add_picture(filename,x,y,w):

    path = FIGURES / filename

    if path.exists():

        slide.shapes.add_picture(
            str(path),
            Inches(x),
            Inches(y),
            width=Inches(w)
        )

    else:

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(4)
        )

        box.fill.background()

        box.line.color.rgb = GRAY

        tf = box.text_frame

        tf.text = filename


def build_qr_image(target_url):

    qr_output = OUTPUT / "poster_qr.png"

    if qrcode is None:
        print("[Warning] 'qrcode' package not available. Using placeholder QR box.")
        return None

    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=12,
        border=2,
    )
    qr.add_data(target_url)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")
    img.save(qr_output)

    return qr_output


def export_pdf(pptx_path, pdf_path):

    if os.name != "nt":
        print("[Info] PDF auto-export is only available on Windows with PowerPoint installed.")
        return False

    if win32com is None:
        print("[Info] pywin32 is unavailable, skipping automatic PDF export.")
        return False

    # PowerPoint format code 32 = PDF
    pp_save_as_pdf = 32

    app = None
    presentation = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1
        presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
        presentation.SaveAs(str(pdf_path), pp_save_as_pdf)
        return True
    except Exception as exc:
        print(f"[Warning] PDF export failed: {exc}")
        return False
    finally:
        if presentation is not None:
            presentation.Close()
        if app is not None:
            app.Quit()

# ----------------------------------------------------------
# Title
# ----------------------------------------------------------

textbox(
    0.4,
    0.2,
    47,
    0.6,
    "Neural Population Dynamics from Internal Representations in Learned Obstacle Avoidance",
    size=30,
    bold=True,
    color=QUEENS_BLUE
)

textbox(
    0.5,
    0.9,
    47,
    0.5,
    "Peter Ohue   •   Emily Oby   •   Gunnar Blohm\n"
    "Centre for Neuroscience Studies, Queen's University",
    size=17
)

# ----------------------------------------------------------
# Logos
# ----------------------------------------------------------

def logo(name,x):

    path = LOGOS / name

    if path.exists():

        slide.shapes.add_picture(
            str(path),
            Inches(x),
            Inches(0.2),
            width=Inches(1.3)
        )

logo("queens_logo.png",43)

logo("connected_minds.png",44.7)

logo("cfref_logo.png",46.2)

# ----------------------------------------------------------
# Column locations
# ----------------------------------------------------------

LEFT = 0.5

MID = 16.3

RIGHT = 32.1

# ----------------------------------------------------------
# Introduction
# ----------------------------------------------------------

heading("Introduction",LEFT,1.6)

textbox(
LEFT,
2.2,
14,
3.8,
"""Artificial intelligence controllers remain difficult to interpret despite
their success in autonomous navigation.

We propose a neuroscience-inspired framework that analyses hidden-layer
activity as neural population recordings.

The objective is to understand how internal representations support
adaptive obstacle avoidance under external perturbations."""
)

# ----------------------------------------------------------
# Contributions
# ----------------------------------------------------------

heading("Key Contributions",LEFT,6.0)

textbox(
LEFT,
6.6,
14,
3.2,
"""• Neuro-inspired analysis of PPO hidden representations

• Low-dimensional neural manifolds

• Behavioural robustness under perturbations

• Distributed neural encoding

• Explainable reinforcement learning"""
)

# ----------------------------------------------------------
# Methods
# ----------------------------------------------------------

heading("Methods",LEFT,10.2)

textbox(
LEFT,
10.8,
14,
2.5,
"""• PPO controller

• Continuous obstacle avoidance

• External perturbations

• PCA

• Neural decoding

• Hidden-unit tuning"""
)

add_picture("figure1.png",0.6,13.2,13)

add_picture("figure2.png",0.6,18.0,13)

# ==========================================================
# PART 2 OF 3
# Results
# ==========================================================

# ----------------------------------------------------------
# RESULTS
# ----------------------------------------------------------

heading("Results", MID, 1.6)

textbox(
    MID,
    2.2,
    14,
    1.2,
    """The trained PPO controller exhibited robust adaptive behaviour under
multiple perturbation conditions. Neural analyses revealed structured latent
representations that evolved continuously during obstacle avoidance.""",
    size=17,
)

# ----------------------------------------------------------
# Behavioural Performance
# ----------------------------------------------------------

heading("Behavioural Performance", MID, 3.8)

textbox(
    MID,
    4.4,
    14,
    1.5,
    """Behavioural metrics demonstrated significant effects of perturbation
magnitude on reward, duration, path length, heading deviation and lateral
error (Kruskal-Wallis, p < 0.001). Maximum speed remained invariant.""",
    size=16,
)

add_picture("figure3.png", 16.5, 6.0, 14)

textbox(
    MID,
    10.1,
    14,
    0.5,
    "Figure 3. Behavioural performance across perturbation conditions.",
    size=13,
    color=GRAY,
)

# ----------------------------------------------------------
# Neural Population Dynamics
# ----------------------------------------------------------

heading("Neural Population Dynamics", MID, 11.0)

textbox(
    MID,
    11.6,
    14,
    1.6,
    """Principal Component Analysis showed that hidden-layer activity evolved
along a structured low-dimensional manifold. Neural trajectories remained
smooth throughout successful goal-directed movement.""",
    size=16,
)

add_picture("figure4.png", 16.5, 13.4, 14)

textbox(
    MID,
    17.5,
    14,
    0.5,
    "Figure 4. Low-dimensional neural manifold revealed by PCA.",
    size=13,
    color=GRAY,
)

# ----------------------------------------------------------
# Hidden Unit Tuning
# ----------------------------------------------------------

heading("Hidden Unit Tuning", MID, 18.2)

textbox(
    MID,
    18.8,
    14,
    1.2,
    """Representative hidden units exhibited diverse tuning relationships with
movement speed, indicating distributed neural encoding rather than
single-unit specialization.""",
    size=16,
)

add_picture("figure5.png", 16.5, 20.2, 14)

textbox(
    MID,
    24.2,
    14,
    0.5,
    "Figure 5. Hidden-unit tuning curves.",
    size=13,
    color=GRAY,
)

# ----------------------------------------------------------
# RIGHT COLUMN
# ----------------------------------------------------------

heading("Interpretability", RIGHT, 1.6)

textbox(
    RIGHT,
    2.2,
    14,
    2.5,
    """Neural population analyses revealed that internal representations
encode behaviourally meaningful information that remains linearly
accessible throughout movement.

These findings support a neuroscience-inspired framework for
interpreting reinforcement learning controllers.""",
    size=17,
)

# ----------------------------------------------------------
# Neural Decoding
# ----------------------------------------------------------

heading("Neural Decoding", RIGHT, 5.2)

textbox(
    RIGHT,
    5.8,
    14,
    1.2,
    """Linear decoders successfully recovered task-relevant variables from
hidden-layer activity, demonstrating that latent representations preserve
behaviourally meaningful information.""",
    size=16,
)

add_picture("figure6.png", 32.3, 7.0, 14)

textbox(
    RIGHT,
    11.2,
    14,
    0.5,
    "Figure 6. Linear decoding of task variables.",
    size=13,
    color=GRAY,
)

# ----------------------------------------------------------
# Decision Making
# ----------------------------------------------------------

heading("Decision Making", RIGHT, 12.0)

textbox(
    RIGHT,
    12.6,
    14,
    1.1,
    """The controller selected alternative trajectories according to the
perturbation condition, demonstrating adaptive route selection instead
of stereotyped behaviour.""",
    size=16,
)

add_picture("figure7.png", 32.3, 13.8, 14)

textbox(
    RIGHT,
    17.8,
    14,
    0.5,
    "Figure 7. Route selection and perturbation statistics.",
    size=13,
    color=GRAY,
)

# ----------------------------------------------------------
# Robustness
# ----------------------------------------------------------

heading("Robustness", RIGHT, 18.6)

textbox(
    RIGHT,
    19.2,
    14,
    1.2,
    """Behaviour remained stable across most perturbation conditions.
Only the strongest perturbations produced substantial degradation in
performance, indicating graceful failure rather than catastrophic collapse.""",
    size=16,
)

add_picture("figure8.png", 32.3, 20.6, 14)

textbox(
    RIGHT,
    24.6,
    14,
    0.5,
    "Figure 8. Robustness across evaluation episodes.",
    size=13,
    color=GRAY,
)

# ==========================================================
# PART 3 OF 3
# Discussion, Conclusion, Footer and Save Poster
# ==========================================================

# ----------------------------------------------------------
# Key Findings
# ----------------------------------------------------------

heading("Key Findings", LEFT, 22.5)

textbox(
    LEFT,
    23.1,
    14,
    2.2,
    """• PPO learned structured low-dimensional neural representations.

• Behaviour changed systematically with perturbation magnitude.

• Hidden units exhibited distributed tuning properties.

• Neural population trajectories evolved smoothly during successful
goal-directed behaviour.

• Internal representations remained interpretable using standard
neuroscience analysis techniques.""",
    size=15,
)

# ----------------------------------------------------------
# Discussion
# ----------------------------------------------------------

heading("Discussion", RIGHT, 25.2)

textbox(
    RIGHT,
    25.8,
    14,
    2.2,
    """Rather than acting as an opaque controller, the trained PPO policy
developed organized internal representations that resemble several
computational principles reported in biological sensorimotor systems.

The findings suggest that neuroscience-inspired analyses can provide
useful insight into how reinforcement learning policies generate
adaptive behaviour.""",
    size=15,
)

# ----------------------------------------------------------
# Conclusion
# ----------------------------------------------------------

heading("Conclusion", MID, 25.2)

textbox(
    MID,
    25.8,
    14,
    1.9,
    """Hidden-layer activity contained meaningful neural population
structure that explained behavioural adaptation under perturbation.

This framework provides a practical approach for improving the
interpretability of reinforcement learning systems operating in
safety-critical environments.""",
    size=15,
)

# ----------------------------------------------------------
# Future Work
# ----------------------------------------------------------

heading("Future Work", LEFT, 25.8)

textbox(
    LEFT,
    26.4,
    14,
    1.8,
    """• Dynamic environments

• Human motor adaptation

• Neuroprosthetic control

• Stroke rehabilitation

• Explainable AI

• Brain-inspired robotics""",
    size=15,
)

# ----------------------------------------------------------
# Footer Background
# ----------------------------------------------------------

footer = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    Inches(0),
    Inches(26.2),
    Inches(48),
    Inches(0.8),
)

footer.fill.solid()
footer.fill.fore_color.rgb = QUEENS_BLUE
footer.line.color.rgb = QUEENS_BLUE

# ----------------------------------------------------------
# Acknowledgements
# ----------------------------------------------------------

textbox(
    0.5,
    26.25,
    18,
    0.45,
    "Acknowledgements",
    size=16,
    bold=True,
    color=WHITE,
)

textbox(
    4.0,
    26.25,
    20,
    0.45,
    "Connected Minds Program • Canada First Research Excellence Fund (CFREF-2022-00010)",
    size=13,
    color=WHITE,
)

# ----------------------------------------------------------
# Contact
# ----------------------------------------------------------

textbox(
    24.0,
    26.25,
    11,
    0.45,
    "Centre for Neuroscience Studies | Queen's University",
    size=13,
    color=WHITE,
)

textbox(
    35.0,
    26.25,
    6,
    0.45,
    "peter.ohue@queensu.ca",
    size=13,
    color=WHITE,
)

# ----------------------------------------------------------
# QR Code Placeholder
# ----------------------------------------------------------

qr_path = LOGOS / "qr_code.png"

if not qr_path.exists():
    generated_qr = build_qr_image(POSTER_QR_URL)
    if generated_qr is not None:
        qr_path = generated_qr

if qr_path.exists():

    slide.shapes.add_picture(
        str(qr_path),
        Inches(46.1),
        Inches(25.9),
        width=Inches(0.75),
    )

else:

    qr = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(46.0),
        Inches(25.95),
        Inches(0.75),
        Inches(0.75),
    )

    qr.fill.background()
    qr.line.color.rgb = WHITE

    tf = qr.text_frame
    tf.text = "QR"

    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.color.rgb = WHITE

# ----------------------------------------------------------
# Connected Minds Logo (Bottom)
# ----------------------------------------------------------

bottom_logos = [
    ("queens_logo.png", 41.0),
    ("connected_minds.png", 42.8),
    ("cfref_logo.png", 44.8),
]

for name, xpos in bottom_logos:

    logo_path = LOGOS / name

    if logo_path.exists():

        slide.shapes.add_picture(
            str(logo_path),
            Inches(xpos),
            Inches(25.95),
            width=Inches(1.2),
        )

# ----------------------------------------------------------
# Save Poster
# ----------------------------------------------------------

output_file = OUTPUT / "Neural_Population_Dynamics_Poster.pptx"
pdf_file = OUTPUT / "Neural_Population_Dynamics_Poster.pdf"

prs.save(output_file)

pdf_exported = export_pdf(output_file, pdf_file)

print("=" * 70)
print("Poster successfully created")
print("=" * 70)
print(f"Editable poster (.pptx): {output_file}")
if pdf_exported:
    print(f"Printable/downloadable poster (.pdf): {pdf_file}")
else:
    print("Printable/downloadable poster (.pdf): not auto-exported")
    print("Open the PPTX in PowerPoint and use File -> Export -> Create PDF/XPS.")
print("=" * 70)